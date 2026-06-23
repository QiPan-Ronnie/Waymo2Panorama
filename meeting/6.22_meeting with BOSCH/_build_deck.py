# -*- coding: utf-8 -*-
"""
Build the Ground-Outpainting deck as a real .pptx, using the old BOSCH/USC deck
as the template (inherits its master/theme/size). 4 slides, drawn programmatically
in the old deck's exact design tokens. Video slots are left as empty placeholders
(the user inserts the mp4s themselves).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

MEET = r"D:\BaiduSyncdisk\2024 to future\koi chen\experiments\Waymo2Panorama\meeting\6.19_meeting with BOSCH"
OLD  = os.path.join(MEET, "Waymo2Pano_BOSCH_progress_2026-06-12.pptx")
OUT  = os.path.join(MEET, "Waymo2Pano_BOSCH_outpainting_2026-06-19.pptx")
ASSET = os.path.join(MEET, "outpainting_deck", "assets")

# ---- design tokens (probed from the old deck) ----
MAROON = RGBColor(0x99,0x00,0x00)   # cardinal: titles, chapter, bars
DEEPRED= RGBColor(0x6E,0x00,0x00)
GOLD   = RGBColor(0xFF,0xCC,0x00)
CREAM  = RGBColor(0xFB,0xF6,0xEC)   # content ground
PAPER  = RGBColor(0xFF,0xFF,0xFF)
INK    = RGBColor(0x1A,0x14,0x10)   # body
DIM    = RGBColor(0x6B,0x60,0x58)   # secondary
LABEL  = RGBColor(0x4A,0x3F,0x33)   # inline label brown
SOFTGOLD = RGBColor(0xFD,0xF3,0xCF) # pale gold fill for highlights/banners
LINE   = RGBColor(0xE3,0xDA,0xC8)   # hairline on cream
# semantic (issues slide)
FIX  = RGBColor(0x2E,0x7D,0x52)
SEMI = RGBColor(0xC9,0x8A,0x00)
PHYS = RGBColor(0x99,0x00,0x00)
ARCH = RGBColor(0x2C,0x52,0x82)
FONT = "Calibri"
EMU_W, EMU_H = 13.333, 7.5

prs = Presentation(OLD)
# wipe existing slides (keep master/theme/layouts) — must drop BOTH the sldId refs
# AND the presentation->slide relationships, else the old slideN.xml parts linger
# in the package and collide with the new ones (duplicate-name corruption).
xml_slides = prs.slides._sldIdLst
for sid in list(xml_slides):
    xml_slides.remove(sid)
_pp = prs.part
for _rid in list(_pp.rels):
    if _pp.rels[_rid].reltype.endswith('/slide'):
        _pp.drop_rel(_rid)
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def new_slide():
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    rect(s, 0, 0, EMU_W, EMU_H, fill=CREAM)            # cream ground
    return s

def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, rounded=False, dash=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
        if dash: shp.line._get_or_add_ln().append(_dash(dash))
    shp.shadow.inherit = False
    if rounded:  # gentle corner radius
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    return shp

def _dash(val):
    from pptx.oxml import parse_xml
    return parse_xml('<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="%s"/>' % val)

def arrow(s, l, t, w, h, color=GOLD):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    try: shp.adjustments[0] = 0.55; shp.adjustments[1] = 0.55
    except Exception: pass
    return shp

def arrow_down(s, l, t, w, h, color=GOLD):
    shp = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def txt(s, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=None):
    """paras = list of paragraphs; each = (align, [ (text,size,bold,color,italic?), ... ])"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for pi, (align, runs) in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        for run in runs:
            tx, sz, b, col = run[0], run[1], run[2], run[3]
            ital = run[4] if len(run) > 4 else False
            r = p.add_run(); r.text = tx
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = col; r.font.italic = ital
    return tb

def spc(run_text_box, pts):
    for p in run_text_box.text_frame.paragraphs:
        for r in p.runs:
            r._r.get_or_add_rPr().set('spc', str(int(pts*100)))

# deck chrome: red/gold top & bottom bars, chapter tag, page numeral
def chrome(s, chapter, numeral):
    rect(s, 0, 0, EMU_W*0.58, 0.085, fill=MAROON)
    rect(s, EMU_W*0.58, 0, EMU_W*0.42, 0.085, fill=GOLD)
    rect(s, 0, EMU_H-0.085, EMU_W*0.42, 0.085, fill=GOLD)
    rect(s, EMU_W*0.42, EMU_H-0.085, EMU_W*0.58, 0.085, fill=MAROON)
    rect(s, 0.62, 0.52, 0.34, 0.045, fill=GOLD)   # chapter rule
    ct = txt(s, 1.04, 0.40, 11.0, 0.32, [(PP_ALIGN.LEFT, [(chapter, 13, True, MAROON)])])
    spc(ct, 2.2)
    txt(s, EMU_W-1.1, EMU_H-0.62, 0.6, 0.34,
        [(PP_ALIGN.RIGHT, [(numeral, 15, True, MAROON)])], anchor=MSO_ANCHOR.MIDDLE)

def title(s, t, em=None):
    runs = [(t, 30, True, MAROON)]
    if em: runs.append((em, 30, True, RGBColor(0x9A,0x6A,0x00), True))
    txt(s, 0.62, 0.80, 12.1, 0.62, [(PP_ALIGN.LEFT, runs)])

# ============================================================
# SLIDE 1 — METHOD: flow diagram
# ============================================================
s = new_slide()
chrome(s, "GROUND OUTPAINTING  ·  METHOD", "I")
txt(s, 0.62, 0.74, 12.1, 0.55, [(PP_ALIGN.LEFT, [
    ("Ground Outpainting", 27, True, MAROON),
    ("  —  Temporal Real-Pixel Reprojection", 27, True, RGBColor(0x9A,0x6A,0x00), True)])])
txt(s, 0.62, 1.33, 12.1, 0.34, [(PP_ALIGN.LEFT, [
    ("Not generative inpainting — real pixels reprojected from other frames.", 13.5, True, RGBColor(0x9A,0x6A,0x00), True)])])

# LEFT: 8-step pipeline (step 4 frame-pick + step 8 abstain highlighted)
steps = [
    ("1","Missing nadir","black hole under / behind the car", False),
    ("2","ERP pixel → 3D ray","each pixel is a viewing ray", False),
    ("3","Ray → ground point","LiDAR-corrected ground height", False),
    ("4","Whole-log frame search","pick by ego displacement, not time", True),
    ("5","Project into source cameras","past / future frames", False),
    ("6","Visibility gates","FOV / ego body / dynamic objects", False),
    ("7","Multi-source consensus","up to 6 sources, median", False),
    ("8","Render or abstain","real pixel, else mask / flat plate", True),
]
lx, lw = 0.62, 5.95
ty, sh, sg = 1.90, 0.485, 0.165
for (n,hd,sb,key) in steps:
    rect(s, lx, ty, lw, sh, fill=(SOFTGOLD if key else PAPER), line=(GOLD if key else LINE), lw=(1.75 if key else 1), rounded=True)
    rect(s, lx+0.11, ty+0.085, 0.33, sh-0.17, fill=(GOLD if key else MAROON), rounded=True)
    txt(s, lx+0.11, ty+0.085, 0.33, sh-0.17, [(PP_ALIGN.CENTER, [(n, 12.5, True, (MAROON if key else PAPER))])], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx+0.56, ty+0.03, lw-0.66, sh-0.06, [
        (PP_ALIGN.LEFT, [(hd, 11, True, MAROON)]),
        (PP_ALIGN.LEFT, [(sb, 9, False, DIM)]),
    ], anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    if n != "8":
        arrow_down(s, lx+lw/2-0.085, ty+sh-0.002, 0.17, sg+0.006)
    ty += sh + sg

# vertical divider
rect(s, 6.86, 1.95, 0.018, 5.0, fill=GOLD)

# RIGHT: explainer (the "why" in plain language; Key idea is the punchline)
rx, rw = 7.10, 5.62
expl = [
    ("Goal", "Fill the missing ground near the vehicle bottom / nadir."),
    ("Key idea", "Use time as the missing camera — the current frame can’t see under the car, but nearby frames observed the same ground point."),
    ("Frame selection", "Not a time window. Search the whole log by ego displacement: ~5–58 m, 5 m buckets, time-nearest in each."),
    ("Rendering", "Project each missing pixel to a 3D ground point, then into candidate source cameras."),
    ("Validity", "Keep only source pixels passing visibility: in FOV, not blocked by ego body or dynamic objects."),
    ("Fusion", "Up to 6 samples → median consensus → render the nearest-to-median real pixel."),
    ("Safety", "Sources disagree or no evidence → abstain (mask / flat plate), never fake texture."),
]
tb = s.shapes.add_textbox(Inches(rx), Inches(1.95), Inches(rw), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
for i,(lab,body) in enumerate(expl):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(8); p.line_spacing = 1.05
    r1 = p.add_run(); r1.text = lab + "   "; r1.font.name = FONT; r1.font.size = Pt(12); r1.font.bold = True
    r1.font.color.rgb = (RGBColor(0x9A,0x6A,0x00) if lab == "Key idea" else MAROON)
    r2 = p.add_run(); r2.text = body; r2.font.name = FONT; r2.font.size = Pt(11); r2.font.bold = False; r2.font.color.rgb = INK

# ============================================================
# SLIDE 2 — RESULT: ground-filled video slots (user inserts mp4s)
# ============================================================
s = new_slide()
chrome(s, "GROUND OUTPAINTING  ·  RESULT", "II")
title(s, "4 scenes · 93 frames each · ", "ground-filled")
txt(s, 0.62, 1.52, 12.1, 0.5, [(PP_ALIGN.LEFT, [
    ("Temporal reprojection — real pixels, no generation. Sky left black (this is the ground layer only).", 14, False, DIM)])])

def video_slots(s, y0, names):
    bw, bh, gap = 4.2, 2.1, 0.34
    total_w = bw*2 + gap
    x0 = (EMU_W - total_w)/2
    pos = [(x0,y0),(x0+bw+gap,y0),(x0,y0+bh+gap),(x0+bw+gap,y0+bh+gap)]
    for (px,py),nm in zip(pos,names):
        rect(s, px, py, bw, bh, fill=RGBColor(0xEE,0xE7,0xD8), line=MAROON, lw=1.25, dash="dash", rounded=True)
        txt(s, px, py+bh/2-0.34, bw, 0.7, [
            (PP_ALIGN.CENTER, [(nm, 16, True, MAROON)]),
            (PP_ALIGN.CENTER, [("▶  insert video", 11, False, DIM)]),
        ], anchor=MSO_ANCHOR.MIDDLE)
    return y0+bh*2+gap

video_slots(s, 2.05, ["bmw","clean","crowd","highway"])
txt(s, 0.62, 6.66, 12.11, 0.4, [(PP_ALIGN.CENTER, [
    ("94.5–100% real-pixel ground   ·   lane lines continuous   ·   ego hood removed   ·   moving cars single & intact", 12.5, True, MAROON)])])

# ============================================================
# SLIDE 3 — OPEN ISSUES + hypotheses
# ============================================================
s = new_slide()
chrome(s, "GROUND OUTPAINTING  ·  OPEN ISSUES", "III")
title(s, "What's wrong — and our ", "hypotheses")

# honest banner
rect(s, 0.62, 1.46, 12.11, 0.82, fill=SOFTGOLD, line=GOLD, lw=1.25, rounded=True)
rect(s, 0.62, 1.46, 0.07, 0.82, fill=GOLD)
txt(s, 0.84, 1.46, 11.85, 0.82, [
    (PP_ALIGN.LEFT, [("These are hypotheses — the data already refuted one of mine.  ", 12, True, MAROON),
                     ("I claimed near-nadir is a universal blind spot; coverage shows most frames have 6-source ~11° cover. It's local, not universal.", 12, False, INK)])],
    anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)

# left: evidence images
imgL, imgW = 0.62, 5.74
cap_h = 0.34
img1 = os.path.join(ASSET, "coverage_highway.png")
img2 = os.path.join(ASSET, "white_blob_highway.png")
y = 2.46
ih = 1.78
if os.path.exists(img1):
    s.shapes.add_picture(img1, Inches(imgL), Inches(y), Inches(imgW), Inches(ih))
txt(s, imgL, y+ih, imgW, cap_h, [(PP_ALIGN.LEFT, [
    ("SOURCE COVERAGE  ", 9.5, True, MAROON),
    ("green = 6 sources · red = few → ground mostly well-seen; only a thin near-car band is starved.", 9.5, False, DIM)])], spacing=1.0)
y2 = y+ih+cap_h+0.16
if os.path.exists(img2):
    s.shapes.add_picture(img2, Inches(imgL), Inches(y2), Inches(imgW), Inches(ih))
txt(s, imgL, y2+ih, imgW, cap_h, [(PP_ALIGN.LEFT, [
    ("highway f85  ", 9.5, True, MAROON),
    ("— “melted nadir / white-blob” + corner melt (annotated).", 9.5, False, DIM)])], spacing=1.0)

# right: 5 hypotheses
hyps = [
    ("①", FIX,  "FIX · MED-HI", "White-blob / soft", " = over-low-pass on well-covered ground, not missing data."),
    ("②", FIX,  "FIX · MED",    "Smear / car-front ghost", " = occlusion not verified. Add LiDAR z-buffer + box gate."),
    ("③", SEMI, "SEMI · MED-HI","Lavender cast (bmw)", " = real grazing Fresnel sky reflection. Only tone-harmonize."),
    ("④", PHYS, "PHYS · MED",   "Blind patches on some frames", " truly unseen (open / low-motion). Local, not universal."),
    ("⑤", ARCH, "ARCH · HIGH",  "Middle band independent & solved", " (off ≡ fill) → Route-2 sidesteps ①②④."),
]
hx, hw = 6.62, 6.11
hy = 2.46; hh = 0.74; hgap = 0.115
for ic, col, conf, head, body in hyps:
    rect(s, hx, hy, hw, hh, fill=PAPER, line=LINE, lw=1, rounded=True)
    rect(s, hx, hy, 0.07, hh, fill=col)
    txt(s, hx+0.18, hy, 0.46, hh, [(PP_ALIGN.LEFT, [(ic, 15, True, col)])], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, hx+0.66, hy+0.05, hw-2.2, hh-0.1, [(PP_ALIGN.LEFT, [
        (head, 11, True, MAROON), (body, 11, False, DIM)])], anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    ct = txt(s, hx+hw-1.5, hy, 1.4, hh, [(PP_ALIGN.RIGHT, [(conf, 8.5, True, col)])], anchor=MSO_ANCHOR.MIDDLE)
    hy += hh + hgap

# verdict line
txt(s, 0.62, 6.96, 12.11, 0.34, [(PP_ALIGN.LEFT, [
    ("VERDICT   ", 11, True, MAROON),
    ("Fixable: white-blob, smear/ghost", 11, True, FIX),
    ("    Semi-physical: lavender", 11, True, SEMI),
    ("    Physical (local): blind patches", 11, True, PHYS)])])

# ============================================================
# SLIDE 4 — ROUTE-2: middle-only (video slots for user)
# ============================================================
s = new_slide()
chrome(s, "ROUTE-2  ·  IN PROGRESS", "IV")
title(s, "Middle-band only — ", "hand ground & sky to the world-model")
txt(s, 0.62, 1.52, 12.1, 0.66, [(PP_ALIGN.LEFT, [
    ("Stitch only what the cameras directly see; nadir & sky left black for the downstream generator (Cosmos) to outpaint under a perfect first frame.", 13.5, False, DIM)])])

# video slots (smaller, leave room for the why-cards)
bw, bh, gap = 3.34, 1.67, 0.30
total_w = bw*2 + gap
x0 = (EMU_W - total_w)/2
y0 = 2.34
pos = [(x0,y0),(x0+bw+gap,y0),(x0,y0+bh+gap),(x0+bw+gap,y0+bh+gap)]
for (px,py),nm in zip(pos,["bmw","clean","crowd","highway"]):
    rect(s, px, py, bw, bh, fill=RGBColor(0xEE,0xE7,0xD8), line=MAROON, lw=1.1, dash="dash", rounded=True)
    txt(s, px, py+bh/2-0.30, bw, 0.62, [
        (PP_ALIGN.CENTER, [(nm, 14, True, MAROON)]),
        (PP_ALIGN.CENTER, [("▶  insert video", 10, False, DIM)]),
    ], anchor=MSO_ANCHOR.MIDDLE)

# why cards (3 across, bottom)
whys = [
    ("No propagated defects", "Skips every ground-fill artifact (white-blob / smear / ghost) by construction."),
    ("What Cosmos wants", "Clean conditioning + honest mask = the masked-360 training distribution."),
    ("First frame perfect", "Frames 2–93 middle-only; nothing fabricated under hard-lock."),
]
wy = 6.02; ww = (12.11 - 2*0.3)/3; wx = 0.62; wh = 0.92
for head, body in whys:
    rect(s, wx, wy, ww, wh, fill=PAPER, line=LINE, lw=1, rounded=True)
    rect(s, wx, wy, 0.07, wh, fill=ARCH)
    txt(s, wx+0.18, wy+0.10, ww-0.32, wh-0.18, [
        (PP_ALIGN.LEFT, [(head, 11.5, True, MAROON)]),
        (PP_ALIGN.LEFT, [(body, 10, False, DIM)]),
    ], spacing=1.0)
    wx += ww + 0.30

txt(s, 0.62, 7.02, 12.11, 0.3, [(PP_ALIGN.LEFT, [
    ("STATUS   ", 10.5, True, MAROON),
    ("rendering on 2× L4 · assemble to mp4 — insert clips into the slots above.", 10.5, False, DIM)])])

prs.save(OUT)
print("SAVED", OUT, "slides=", len(prs.slides._sldIdLst))
