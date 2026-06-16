# -*- coding: utf-8 -*-
"""Build the BOSCH PROGRESS-REPORT deck (7 slides) from the editorial template.

Separate output from the tech deck (Waymo2Pano_BOSCH_2026-06-12.pptx) — this writes
Waymo2Pano_BOSCH_progress_2026-06-12.pptx and never touches the tech version.
"""
import os
import time
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
W2P = r"D:\BaiduSyncdisk\2024 to future\koi chen\experiments\Waymo2Panorama"
SRC = os.path.join(HERE, "neurogen-poisoning-editable.pptx")
OUT = os.path.join(HERE, "Waymo2Pano_BOSCH_progress_2026-06-12.pptx")

# ---- palette (matches template) ----
MAROON      = RGBColor(0x6E, 0x00, 0x00)   # title bg
DARKRED     = RGBColor(0x99, 0x00, 0x00)   # eyebrow / title / accents
GOLD        = RGBColor(0xFF, 0xCC, 0x00)   # accent
SOFTGOLD    = RGBColor(0xF4, 0xE4, 0xA1)   # soft gold
CREAM       = RGBColor(0xFB, 0xF6, 0xEC)   # content bg / light text
INK         = RGBColor(0x1A, 0x14, 0x10)   # body text
GRAY_MUTED  = RGBColor(0x6B, 0x60, 0x58)   # secondary text
CAPTION     = RGBColor(0x4A, 0x3F, 0x33)   # caption text
CARD        = RGBColor(0xF1, 0xE9, 0xDA)   # subtle card panel on cream
CARDLINE    = RGBColor(0xCC, 0xC0, 0xA8)   # faint rule
FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)
ROMAN = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X", "XI", "XII"]


def _set_font(run, size=None, bold=None, color=None, italic=None, name=FONT):
    f = run.font
    f.name = name
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color


def _no_autofit(tf):
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, l, t, w, h, runs_or_text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             space_after=2, wrap=True, italic=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs_or_text, str):
        paras = [[(runs_or_text, {})]]
    else:
        paras = runs_or_text
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        for (txt, ov) in para:
            r = p.add_run()
            r.text = txt
            _set_font(r,
                      size=ov.get('size', size),
                      bold=ov.get('bold', bold),
                      color=ov.get('color', color),
                      italic=ov.get('italic', italic),
                      name=ov.get('name', FONT))
    return tb


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE, shadow_off=True):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w if line_w else Pt(1)
    if shadow_off:
        sp.shadow.inherit = False
    return sp


def clear_slide(slide):
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def base_content(slide, idx, eyebrow, title, title_size=30):
    """Editorial chrome for a content slide (cream bg)."""
    add_rect(slide, 0, 0, SW, SH, fill=CREAM)
    # top split band (dark red | gold)
    add_rect(slide, 0, 0, Inches(8.0), Inches(0.07), fill=DARKRED)
    add_rect(slide, Inches(8.0), 0, Inches(5.33), Inches(0.07), fill=GOLD)
    # bottom split band (gold | dark red)
    add_rect(slide, 0, Inches(7.46), Inches(4.7), Inches(0.04), fill=GOLD)
    add_rect(slide, Inches(4.7), Inches(7.46), Inches(8.63), Inches(0.04), fill=DARKRED)
    # roman numeral bottom-right
    add_text(slide, Inches(11.73), Inches(6.95), Inches(1.1), Inches(0.4),
             ROMAN[idx], size=14, bold=False, color=DARKRED, align=PP_ALIGN.RIGHT)
    # eyebrow tick + word
    add_rect(slide, Inches(0.7), Inches(0.62), Inches(0.45), Inches(0.03), fill=GOLD)
    add_text(slide, Inches(1.3), Inches(0.5), Inches(11.0), Inches(0.38),
             eyebrow, size=13, bold=True, color=DARKRED, align=PP_ALIGN.LEFT)
    # big title
    add_text(slide, Inches(0.7), Inches(0.92), Inches(11.93), Inches(1.0),
             title, size=title_size, bold=True, color=DARKRED, align=PP_ALIGN.LEFT,
             line_spacing=1.0)


def fit_image(path, box_l, box_t, box_w, box_h, align='center', valign='middle'):
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    bw, bh = box_w, box_h
    bar = bw / bh
    if ar > bar:
        w = bw; h = int(bw / ar)
    else:
        h = bh; w = int(bh * ar)
    if align == 'center':
        l = box_l + (bw - w) // 2
    elif align == 'left':
        l = box_l
    else:
        l = box_l + (bw - w)
    if valign == 'middle':
        t = box_t + (bh - h) // 2
    elif valign == 'top':
        t = box_t
    else:
        t = box_t + (bh - h)
    return l, t, w, h


def add_image(slide, path, box_l, box_t, box_w, box_h, align='center',
              valign='middle', frame=True):
    l, t, w, h = fit_image(path, box_l, box_t, box_w, box_h, align, valign)
    pic = slide.shapes.add_picture(path, l, t, w, h)
    if frame:
        pic.line.color.rgb = CARDLINE
        pic.line.width = Pt(1)
    return (round(l / 914400, 2), round(t / 914400, 2),
            round(w / 914400, 2), round(h / 914400, 2))


def add_image_exact(slide, path, l, w, t, frame=True):
    """Place image at exact left/top/width, height derived from aspect."""
    im = Image.open(path)
    iw, ih = im.size
    h = int(w / (iw / ih))
    pic = slide.shapes.add_picture(path, l, t, w, h)
    if frame:
        pic.line.color.rgb = CARDLINE
        pic.line.width = Pt(1)
    return (round(l / 914400, 2), round(t / 914400, 2),
            round(w / 914400, 2), round(h / 914400, 2))


def bullets(items, size=14, color=INK, lead_color=GOLD):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            txt, ov = it
        else:
            txt, ov = it, {}
        paras.append([("▸  ", {'color': lead_color, 'bold': True, 'size': size}),
                      (txt, {'size': ov.get('size', size),
                             'color': ov.get('color', color),
                             'bold': ov.get('bold', False),
                             'italic': ov.get('italic', None)})])
    return paras


P = os.path.join


def img(rel):
    return P(W2P, rel)


GEOM = {}

# =====================================================================
prs = Presentation(SRC)
slides = prs.slides

# ---------------------------------------------------------------- S1 (I  title)
s = slides[0]
clear_slide(s)
add_rect(s, 0, 0, SW, SH, fill=MAROON)
# top bands gold | cream
add_rect(s, 0, 0, Inches(8.0), Inches(0.07), fill=GOLD)
add_rect(s, Inches(8.0), 0, Inches(5.33), Inches(0.07), fill=CREAM)
# bottom bands cream | gold
add_rect(s, 0, Inches(7.43), Inches(4.7), Inches(0.05), fill=CREAM)
add_rect(s, Inches(4.7), Inches(7.43), Inches(8.63), Inches(0.05), fill=GOLD)
# roman
add_text(s, Inches(11.73), Inches(6.95), Inches(1.1), Inches(0.4), "I",
         size=14, color=GOLD, align=PP_ALIGN.RIGHT)
# eyebrow
add_text(s, 0, Inches(1.15), SW, Inches(0.4),
         "BOSCH RESEARCH SYNC  ·  2026-06-12", size=14, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
# big title (two lines)
add_text(s, Inches(0.7), Inches(2.0), Inches(11.93), Inches(1.9),
         [[("Perspective Images → 360° Panorama", {})],
          [("Progress Update", {'size': 30, 'color': CREAM})]],
         size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER, line_spacing=1.05)
# divider diamond
add_rect(s, Inches(5.55), Inches(4.18), Inches(0.9), Inches(0.012), fill=GOLD)
add_rect(s, Inches(6.6), Inches(4.11), Inches(0.13), Inches(0.13), fill=GOLD,
         shape=MSO_SHAPE.DIAMOND)
add_rect(s, Inches(6.85), Inches(4.18), Inches(0.9), Inches(0.012), fill=GOLD)
# subtitle
add_text(s, Inches(1.4), Inches(4.5), Inches(10.53), Inches(0.9),
         "From 7 ring-camera images to a complete-sphere ERP panorama (1024×2048).",
         size=18, color=CREAM, align=PP_ALIGN.CENTER, line_spacing=1.15)
# presenter chip
add_rect(s, Inches(5.57), Inches(5.7), Inches(2.2), Inches(0.55), fill=None,
         line=GOLD, line_w=Pt(1.25))
add_text(s, Inches(5.57), Inches(5.79), Inches(2.2), Inches(0.4),
         "JINGSHUO", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(6.42), SW, Inches(0.4),
         "Presenter", size=12, color=CREAM, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S2 (II  DATASET)
s = slides[1]
clear_slide(s)
base_content(s, 1, "DATASET", "Dataset: Argoverse 2 Sensor", title_size=30)
# three stat callouts
def stat(slide, lx, num, lab):
    add_rect(slide, Inches(lx), Inches(2.35), Inches(0.04), Inches(1.4), fill=GOLD)
    add_text(slide, Inches(lx + 0.22), Inches(2.3), Inches(3.4), Inches(0.95),
             num, size=42, bold=True, color=DARKRED)
    add_text(slide, Inches(lx + 0.22), Inches(3.32), Inches(3.4), Inches(0.55),
             lab, size=11.5, bold=True, color=CAPTION)
stat(s, 0.7, "7", "RING CAMERAS · FULL 360°")
stat(s, 4.7, "20 Hz", "CAPTURE RATE")
stat(s, 8.7, "~300", "FRAMES PER LOG")
# bullets below (explicit paragraph lists so the 3rd item can carry inline emphasis)
LEAD = ("▸  ", {'color': GOLD, 'bold': True, 'size': 16})
add_text(s, Inches(0.7), Inches(4.25), Inches(11.93), Inches(2.6),
         [
           [LEAD,
            ("7 ring cameras give full 360° horizontal coverage · 20 Hz · ~300 frames per log.",
             {'size': 16})],
           [LEAD,
            ("3D box annotations + LiDAR + per-camera capture timestamps for every frame.",
             {'size': 16})],
           [LEAD,
            ("Waymo Open has 5 cameras (~250° forward, rear gap) — ", {'size': 16}),
            ("no full ring from cameras alone", {'bold': True, 'color': DARKRED, 'size': 16}),
            ("; Waymo migration is queued as our generalization gate.", {'size': 16})],
         ], line_spacing=1.2, space_after=14)

# ---------------------------------------------------------------- S3 (III CHALLENGES)
s = slides[2]
clear_slide(s)
base_content(s, 2, "CHALLENGES", "Stitching perspective views: what goes wrong",
             title_size=28)
# three stacked problem blocks on the left, shutter diagram on the right
card_l = Inches(0.7)
card_w = Inches(6.75)
card_h = Inches(1.4)
card_tops = [2.25, 3.83, 5.41]
cards = [
    ("①", "Different positions",
     "Each camera sees the world from a different spot — structures break at seams, overlaps ghost."),
    ("②", "Different instants",
     "The 7 shutters fire up to ±22.5 ms apart — a moving car is photographed in two places (a TIME problem; no alignment fixes it)."),
    ("③", "Missing coverage",
     "Sky above and ground below sit outside every camera's FOV — black holes."),
]
for (badge, head, body), top in zip(cards, card_tops):
    add_rect(s, card_l, Inches(top), card_w, card_h, fill=CARD)
    add_rect(s, card_l, Inches(top), Inches(0.06), card_h, fill=DARKRED)
    add_text(s, card_l + Inches(0.22), Inches(top + 0.13), Inches(0.7), Inches(0.7),
             badge, size=26, bold=True, color=GOLD)
    add_text(s, card_l + Inches(0.85), Inches(top + 0.13), card_w - Inches(1.0), Inches(0.45),
             head, size=16, bold=True, color=DARKRED)
    add_text(s, card_l + Inches(0.85), Inches(top + 0.58), card_w - Inches(1.05), Inches(0.78),
             body, size=12.5, color=INK, line_spacing=1.08)
# shutter diagram on the right
GEOM['D2_shutter'] = add_image(s, img(r"meeting\6.12_meeting with BOSCH\assets\_trimmed\D2_shutter.png"),
                               Inches(7.75), Inches(2.55), Inches(5.0), Inches(3.6),
                               align='center', valign='middle')
add_text(s, Inches(7.75), Inches(6.25), Inches(5.0), Inches(0.5),
         "7 rolling shutters across one frame interval.",
         size=11, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S4 (IV METHOD)
s = slides[3]
clear_slide(s)
base_content(s, 3, "OUR METHOD", "Our method: remove the cause — three fixes",
             title_size=28)
blocks = [
    [("①  Virtual centre  ", {'bold': True, 'color': DARKRED, 'size': 15}),
     ("— project from the cameras' own centroid instead of the ego origin → projection residuals ↓ 18–96×, seams align.",
      {'size': 14})],
    [("②  EMC per-camera shutter poses  ", {'bold': True, 'color': DARKRED, 'size': 15}),
     ("— each image projected with the ego pose of its own exposure instant (SE(3) interpolation) → the static world aligns exactly.",
      {'size': 14})],
    [("③  One body · one camera · one time (+OMC)  ", {'bold': True, 'color': DARKRED, 'size': 15}),
     ("— every moving object rendered whole from a single camera at a single instant (YOLOv8 masks + 3D-box identity, ECC-measured shutter shift) → ghosted / seam-cut cars impossible by construction.",
      {'size': 14})],
]
btops = [2.05, 2.93, 3.81]
for blk, top in zip(blocks, btops):
    add_rect(s, Inches(0.7), Inches(top), Inches(0.04), Inches(0.74), fill=GOLD)
    add_text(s, Inches(0.92), Inches(top - 0.02), Inches(11.7), Inches(0.82),
             [blk], line_spacing=1.06, space_after=0)
# result line with gold tick
add_rect(s, Inches(0.7), Inches(4.68), Inches(11.93), Inches(0.46), fill=CARD)
add_text(s, Inches(0.9), Inches(4.74), Inches(11.5), Inches(0.4),
         [[("✓  ", {'color': GOLD, 'bold': True, 'size': 15}),
           ("Zero per-scene parameters", {'bold': True, 'color': DARKRED, 'size': 13.5}),
           (" — the same code passes 5 scenes (bmw / downtown / crowd / clean / highway), graceful without LiDAR.",
            {'size': 13.5, 'color': INK})]],
         anchor=MSO_ANCHOR.MIDDLE)
# result image (scene band) — sized so its bottom stays >=0.35in clear of the band.
# AR ~3.03 is very wide/short; the 3 blocks + result line leave <2.6in of height,
# so the image is centred at the width that physically fits the remaining band.
_fig_w = 5.4
GEOM['fig2_bmw'] = add_image_exact(
    s, img(r"agent\figs\fig2_bmw_final_pano.jpg"),
    l=Inches((13.333 - _fig_w) / 2), w=Inches(_fig_w), t=Inches(5.32))

# ---------------------------------------------------------------- S5 (V SPHERE)
s = slides[4]
clear_slide(s)
base_content(s, 4, "COMPLETING THE SPHERE", "Outpainting the missing sphere: ground + sky",
             title_size=27)
col_w = Inches(5.85)
# LEFT column - GROUND
lx = 0.7
add_text(s, Inches(lx), Inches(2.15), col_w, Inches(0.4),
         [[("GROUND", {'bold': True, 'color': DARKRED, 'size': 15}),
           ("  — real pixels, no generation", {'color': GRAY_MUTED, 'size': 13})]])
add_rect(s, Inches(lx), Inches(2.6), col_w, Inches(0.014), fill=GOLD)
add_text(s, Inches(lx), Inches(2.78), col_w, Inches(3.9),
         [
           [("INPUT  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("scene-band panorama + the log's other frames (ego moved).", {'size': 13.5})],
           [("Method  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("temporal reprojection — the road under the car now was photographed seconds earlier / later; reproject those real pixels (occlusion-checked).",
             {'size': 13.5})],
           [("OUTPUT  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("nadir filled with 94.5–100% real pixels, lane lines continuous, ego hood removed.",
             {'size': 13.5})],
         ], line_spacing=1.15, space_after=10)
# RIGHT column - SKY
rx = 7.0
add_text(s, Inches(rx), Inches(2.15), col_w, Inches(0.4),
         [[("SKY", {'bold': True, 'color': DARKRED, 'size': 15}),
           ("  — the only generated layer", {'color': GRAY_MUTED, 'size': 13})]])
add_rect(s, Inches(rx), Inches(2.6), col_w, Inches(0.014), fill=GOLD)
add_text(s, Inches(rx), Inches(2.78), col_w, Inches(3.9),
         [
           [("INPUT  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("panorama + sky mask (the only evidence-free region).", {'size': 13.5})],
           [("Method  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("FLUX.1-Fill diffusion inpainting + DiT360 panorama LoRA, auto-prompt from the visible sky band (5/5 scenes correct).",
             {'size': 13.5})],
           [("OUTPUT  ", {'bold': True, 'color': DARKRED, 'size': 13}),
            ("complete sphere — every pixel outside the mask byte-identical to input.",
             {'size': 13.5})],
         ], line_spacing=1.15, space_after=10)
# I/O diagram full width below
GEOM['D4_outpaint'] = add_image_exact(
    s, img(r"meeting\6.12_meeting with BOSCH\assets\_trimmed\D4_outpaint.png"),
    l=Inches((13.333 - 6.5) / 2), w=Inches(6.5), t=Inches(5.05))

# ---------------------------------------------------------------- S6 (VI RESULTS)
s = slides[5]
clear_slide(s)
base_content(s, 5, "RESULTS", "Current results: complete spheres, multi-weather",
             title_size=27)
# hero image (W 6.6in centered, top ~1.85)
GEOM['downtown'] = add_image_exact(
    s, img(r"deliverables\complete_pano_v8\downtown_complete_v8.png"),
    l=Inches((13.333 - 6.6) / 2), w=Inches(6.6), t=Inches(1.85))
# two thumbnails W 3.22 each, gap 0.16, below hero
thumb_w = 3.22
gap = 0.16
total = thumb_w * 2 + gap
tx0 = (13.333 - total) / 2
ty = 5.15
GEOM['bmw'] = add_image_exact(
    s, img(r"deliverables\complete_pano_v8\bmw_complete_v8.png"),
    l=Inches(tx0), w=Inches(thumb_w), t=Inches(ty))
GEOM['highway'] = add_image_exact(
    s, img(r"deliverables\complete_pano_v8\highway_complete_v8.png"),
    l=Inches(tx0 + thumb_w + gap), w=Inches(thumb_w), t=Inches(ty))
# caption
add_text(s, Inches(0.7), Inches(7.0), Inches(11.93), Inches(0.42),
         "downtown auto-detected dusk · bmw sunny · highway altocumulus — "
         "scene band & ground 100% real pixels; sky is the only generated region.",
         size=11, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S7 (VII NEXT)
s = slides[6]
clear_slide(s)
base_content(s, 6, "NEXT", "Next steps", title_size=30)
add_text(s, Inches(0.7), Inches(2.5), Inches(11.93), Inches(3.5),
         [
           [("▸  ", {'color': GOLD, 'bold': True, 'size': 17}),
            ("Waymo migration — the generalization gate", {'bold': True, 'color': DARKRED, 'size': 17}),
            ("  (loader-level changes only?).", {'size': 17})],
           [("▸  ", {'color': GOLD, 'bold': True, 'size': 17}),
            ("Centre / format contract", {'bold': True, 'color': DARKRED, 'size': 17}),
            ("  with the world-model first-frame consumer.", {'size': 17})],
           [("▸  ", {'color': GOLD, 'bold': True, 'size': 17}),
            ("Scale", {'bold': True, 'color': DARKRED, 'size': 17}),
            ("  — a 75-panorama AV2 set already produced; extend with the v8 complete-sphere pipeline.",
             {'size': 17})],
         ], line_spacing=1.2, space_after=20)

# ---------------------------------------------------------------- delete surplus 8..13
xml_slides = prs.slides._sldIdLst
sldIds = list(xml_slides)
for sid in sldIds[7:]:
    xml_slides.remove(sid)

# ---------------------------------------------------------------- save with retry
saved = False
for attempt in range(6):
    try:
        prs.save(OUT)
        saved = True
        break
    except Exception as e:
        print(f"save attempt {attempt + 1} failed: {e}; retrying in 5s")
        time.sleep(5)
if not saved:
    raise SystemExit("ERROR: could not save after 6 attempts")

print("SAVED", OUT)
print("slides:", len(prs.slides))
print("IMAGE GEOM (l,t,w,h inches):")
for k, v in GEOM.items():
    print(f"  {k}: {v}")
