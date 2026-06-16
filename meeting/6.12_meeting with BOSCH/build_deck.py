# -*- coding: utf-8 -*-
"""Build the BOSCH research-sync deck from the editorial template."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
W2P = r"D:\BaiduSyncdisk\2024 to future\koi chen\experiments\Waymo2Panorama"
SRC = os.path.join(HERE, "neurogen-poisoning-editable.pptx")
OUT = os.path.join(HERE, "Waymo2Pano_BOSCH_2026-06-12.pptx")

# ---- palette ----
MAROON      = RGBColor(0x6E, 0x00, 0x00)   # title bg
DARKRED     = RGBColor(0x99, 0x00, 0x00)   # eyebrow / title / accents
GOLD        = RGBColor(0xFF, 0xCC, 0x00)   # accent
CREAM       = RGBColor(0xFB, 0xF6, 0xEC)   # content bg / light text
INK         = RGBColor(0x1A, 0x14, 0x10)   # body text
GRAY_MUTED  = RGBColor(0x6B, 0x60, 0x58)   # secondary text
CARD        = RGBColor(0xF1, 0xE9, 0xDA)   # subtle card panel on cream
PLACE_FILL  = RGBColor(0xE4, 0xDD, 0xD0)   # placeholder light gray
PLACE_LINE  = RGBColor(0xB9, 0xAE, 0x9C)
FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

ROMAN = ["I","II","III","IV","V","VI","VII","VIII","IX","X","XI","XII"]


def _set_font(run, size=None, bold=None, color=None, italic=None, name=FONT):
    f = run.font
    f.name = name
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color


def _no_autofit(tf):
    # disable autosize so text keeps our point sizes
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, l, t, w, h, runs_or_text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             space_after=2, wrap=True):
    """runs_or_text: str, or list of paragraphs where each paragraph is a list
    of (text, dict-overrides) run tuples."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs_or_text, str):
        paras = [[(runs_or_text, {})]]
    else:
        paras = runs_or_text

    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        for (txt, ov) in para:
            r = p.add_run()
            r.text = txt
            _set_font(r,
                      size=ov.get('size', size),
                      bold=ov.get('bold', bold),
                      color=ov.get('color', color),
                      italic=ov.get('italic', None),
                      name=ov.get('name', FONT))
    return tb


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE, shadow_off=True):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w if line_w else Pt(1)
    if shadow_off:
        sp.shadow.inherit = False
    return sp


def clear_slide(slide):
    """Remove all shapes from a slide."""
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def base_content(slide, idx, eyebrow, title, title_size=30):
    """Lay down the editorial chrome for a content slide (cream bg)."""
    # full-bleed cream background
    add_rect(slide, 0, 0, SW, SH, fill=CREAM)
    # top split band (dark red | gold)
    add_rect(slide, 0, 0, Inches(8.0), Inches(0.07), fill=DARKRED)
    add_rect(slide, Inches(8.0), 0, Inches(5.33), Inches(0.07), fill=GOLD)
    # bottom split band (dark red | gold)
    add_rect(slide, 0, Inches(7.46), Inches(4.7), Inches(0.04), fill=DARKRED)
    add_rect(slide, Inches(4.7), Inches(7.46), Inches(8.63), Inches(0.04), fill=GOLD)
    # roman numeral bottom-right
    add_text(slide, Inches(11.73), Inches(6.95), Inches(1.1), Inches(0.4),
             ROMAN[idx], size=14, bold=False, color=DARKRED, align=PP_ALIGN.RIGHT)
    # eyebrow tick + word
    add_rect(slide, Inches(0.7), Inches(0.62), Inches(0.45), Inches(0.03), fill=GOLD)
    add_text(slide, Inches(1.3), Inches(0.5), Inches(9.5), Inches(0.38),
             eyebrow, size=13, bold=True, color=DARKRED, align=PP_ALIGN.LEFT)
    # big title
    add_text(slide, Inches(0.7), Inches(0.92), Inches(11.93), Inches(1.0),
             title, size=title_size, bold=True, color=DARKRED, align=PP_ALIGN.LEFT,
             line_spacing=1.0)
    # accent tick under title
    add_rect(slide, Inches(0.7), Inches(1.92), Inches(0.7), Inches(0.04), fill=GOLD)


def placeholder(slide, tag, l, t, w, h):
    """Light-gray rounded rect with centered [[Dx]] text."""
    sp = add_rect(slide, l, t, w, h, fill=PLACE_FILL, line=PLACE_LINE,
                  line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # softer corner radius
    try:
        sp.adjustments[0] = 0.04
    except Exception:
        pass
    tf = sp.text_frame
    tf.word_wrap = True
    _no_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    _set_font(r, size=24, bold=True, color=GRAY_MUTED)
    return (round(l/914400, 2), round(t/914400, 2), round(w/914400, 2), round(h/914400, 2))


def fit_image(path, box_l, box_t, box_w, box_h, align='center', valign='middle'):
    """Return (l,t,w,h) EMU fitting image into box keeping aspect ratio."""
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    bw, bh = box_w, box_h
    bar = bw / bh
    if ar > bar:        # image wider -> fit width
        w = bw; h = int(bw / ar)
    else:               # fit height
        h = bh; w = int(bh * ar)
    if align == 'center':
        l = box_l + (bw - w) // 2
    elif align == 'left':
        l = box_l
    else:
        l = box_l + (bw - w)
    if valign == 'middle':
        t = box_t + (bh - h) // 2
    elif valign == 'top':
        t = box_t
    else:
        t = box_t + (bh - h)
    return l, t, w, h


def add_image(slide, path, box_l, box_t, box_w, box_h, align='center', valign='middle'):
    l, t, w, h = fit_image(path, box_l, box_t, box_w, box_h, align, valign)
    slide.shapes.add_picture(path, l, t, w, h)
    return (round(l/914400, 2), round(t/914400, 2), round(w/914400, 2), round(h/914400, 2))


# bullet helper: build paragraph list with a leading dot
def bullets(items, size=13, color=INK, lead_color=GOLD, sep_after=5):
    paras = []
    for it in items:
        if isinstance(it, tuple):  # (text, override-dict)
            txt, ov = it
        else:
            txt, ov = it, {}
        paras.append([("▸  ", {'color': lead_color, 'bold': True, 'size': size}),
                      (txt, {'size': ov.get('size', size), 'color': ov.get('color', color),
                             'bold': ov.get('bold', False)})])
    return paras


P = os.path.join

def img(rel):
    return P(W2P, rel)

# track placeholder geometries to report
PLACE_GEOM = {}

# =====================================================================
prs = Presentation(SRC)
slides = prs.slides

# We have 13 slides; we will rebuild slides 0..11 (12 slides) and delete slide 12.

# ---------------------------------------------------------------- S1 (title)
s = slides[0]
clear_slide(s)
add_rect(s, 0, 0, SW, SH, fill=MAROON)
# top bands gold|cream
add_rect(s, 0, 0, Inches(8.0), Inches(0.07), fill=GOLD)
add_rect(s, Inches(8.0), 0, Inches(5.33), Inches(0.07), fill=CREAM)
# bottom bands cream|gold
add_rect(s, 0, Inches(7.43), Inches(4.7), Inches(0.05), fill=CREAM)
add_rect(s, Inches(4.7), Inches(7.43), Inches(8.63), Inches(0.05), fill=GOLD)
# roman
add_text(s, Inches(11.73), Inches(6.95), Inches(1.1), Inches(0.4), "I",
         size=14, color=GOLD, align=PP_ALIGN.RIGHT)
# eyebrow
add_text(s, 0, Inches(1.05), SW, Inches(0.4),
         "BOSCH RESEARCH SYNC  ·  2026-06-12", size=14, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
# big title
add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(1.5),
         "Perspective Images → 360° Panorama", size=52, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
# divider diamond
add_rect(s, Inches(5.7), Inches(3.62), Inches(0.9), Inches(0.012), fill=GOLD)
add_rect(s, Inches(6.65), Inches(3.55), Inches(0.13), Inches(0.13), fill=GOLD,
         shape=MSO_SHAPE.DIAMOND)
add_rect(s, Inches(6.85), Inches(3.62), Inches(0.9), Inches(0.012), fill=GOLD)
# subtitle
add_text(s, Inches(1.4), Inches(3.95), Inches(10.53), Inches(1.1),
         "A general multi-camera pipeline from AV ring cameras to a complete-sphere "
         "ERP panorama — source-faithful, evidence-gated.",
         size=17, color=CREAM, align=PP_ALIGN.CENTER, line_spacing=1.15)
# presenter chip
add_rect(s, Inches(5.57), Inches(5.55), Inches(2.2), Inches(0.55), fill=None,
         line=GOLD, line_w=Pt(1.25))
add_text(s, Inches(5.57), Inches(5.64), Inches(2.2), Inches(0.4),
         "JINGSHUO", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(6.35), SW, Inches(0.4),
         "Presenter", size=12, color=CREAM, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- S2 (II OPENING)
s = slides[1]
clear_slide(s)
base_content(s, 1, "OPENING", "The task — and why Argoverse 2", title_size=30)
# left goal column
add_text(s, Inches(0.7), Inches(2.25), Inches(5.55), Inches(0.34),
         "GOAL", size=12, bold=True, color=GRAY_MUTED)
add_rect(s, Inches(0.7), Inches(2.62), Inches(5.55), Inches(0.012), fill=GOLD)
add_text(s, Inches(0.7), Inches(2.78), Inches(5.55), Inches(2.2),
         bullets([
           "N synchronized perspective cameras → one 1024×2048 ERP 360° panorama",
           "Downstream consumer: world-model (Cosmos-style) first frame — provenance must be honest",
         ], size=14), color=INK, line_spacing=1.12, space_after=10)
# right comparison: two cards
card_t = Inches(2.25)
card_h = Inches(3.95)
# Waymo card
add_rect(s, Inches(6.55), card_t, Inches(2.95), card_h, fill=CARD)
add_rect(s, Inches(6.55), card_t, Inches(0.06), card_h, fill=GRAY_MUTED)
add_text(s, Inches(6.75), Inches(2.4), Inches(2.6), Inches(0.3),
         "WAYMO OPEN DATASET", size=12, bold=True, color=DARKRED)
add_text(s, Inches(6.75), Inches(2.78), Inches(2.6), Inches(3.3),
         bullets([
           "5 cameras",
           "~250° forward arc",
           "rear gap",
           ("No full ring from cameras alone", {'bold': True, 'color': DARKRED}),
         ], size=12.5), line_spacing=1.1, space_after=7)
# AV2 card (highlighted)
add_rect(s, Inches(9.65), card_t, Inches(3.0), card_h, fill=CARD)
add_rect(s, Inches(9.65), card_t, Inches(0.06), card_h, fill=GOLD)
add_text(s, Inches(9.85), Inches(2.4), Inches(2.65), Inches(0.3),
         "ARGOVERSE 2", size=12, bold=True, color=DARKRED)
add_text(s, Inches(9.85), Inches(2.78), Inches(2.65), Inches(3.3),
         bullets([
           "7 ring cameras = full 360°",
           "20 Hz · ~300 frames/log",
           "3D boxes + LiDAR",
           ("Per-camera capture timestamps (this became crucial)", {'bold': True, 'color': DARKRED}),
         ], size=12.5), line_spacing=1.1, space_after=7)
# footer line
add_rect(s, Inches(0.7), Inches(6.45), Inches(0.04), Inches(0.55), fill=GOLD)
add_text(s, Inches(0.9), Inches(6.5), Inches(11.4), Inches(0.5),
         "Waymo migration stays queued as our generalization gate.",
         size=13, bold=True, color=DARKRED)

# ---------------------------------------------------------------- S3 (III BASELINE)
s = slides[2]
clear_slide(s)
base_content(s, 2, "BASELINE", "L1 baseline: project and blend — and why it fails",
             title_size=27)
# image (wide 3-panel) across top
g = add_image(s, img(r"deliverables\ghostkill\GK_bmw_avg_vs_pick.png"),
              Inches(0.7), Inches(2.15), Inches(11.93), Inches(2.95),
              align='center', valign='top')
# bullets below
add_text(s, Inches(0.7), Inches(5.35), Inches(11.93), Inches(1.6),
         bullets([
           "Spherical projection from the ego origin + multi-band blending",
           "Symptoms: ghosting in every overlap, structural misalignment",
           ("Diagnosis: blending = averaging two misaligned copies of the world",
            {'bold': True, 'color': DARKRED}),
         ], size=14), line_spacing=1.1, space_after=8)

# ---------------------------------------------------------------- S4 (IV ITERATIONS)
s = slides[3]
clear_slide(s)
base_content(s, 3, "ITERATIONS", "Two months of better seam-hiding", title_size=30)
thumbs = [
    (r"deliverables\gpt_pro_sources\04_L1_hard_select_bmw_2048x1024.png",
     "Hard select",
     "One camera per pixel: color ghosts gone, structure still misaligned"),
    (r"deliverables\gpt_pro_sources\01_A1_view_none_bmw_2048x1024.png",
     "Optical-flow view interpolation",
     "Smooths determinable seams, helpless on moving objects"),
    (r"deliverables\gpt_pro_sources\02_G_bmw_pano_2048x1024.jpg",
     "Seam routing around objects",
     "Best of this era — near-ground waviness and seam-cut cars persist"),
]
col_w = Inches(3.85)
gap = Inches(0.19)
x0 = Inches(0.7)
for i, (rel, head, verdict) in enumerate(thumbs):
    cx = x0 + i * (col_w + gap)
    add_image(s, img(rel), cx, Inches(2.3), col_w, Inches(2.0),
              align='center', valign='top')
    add_text(s, cx, Inches(4.4), col_w, Inches(0.3), head, size=13, bold=True,
             color=DARKRED)
    add_text(s, cx, Inches(4.72), col_w, Inches(1.1), verdict, size=12,
             color=INK, line_spacing=1.08)
# punchline bar
add_rect(s, Inches(0.7), Inches(6.0), Inches(11.93), Inches(0.95), fill=DARKRED)
add_text(s, Inches(0.95), Inches(6.12), Inches(11.4), Inches(0.75),
         "Every route hit the same wall: we were hiding seams, not removing their "
         "cause. Learned / 3D replacements tested worse.",
         size=14, bold=True, color=CREAM, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# ---------------------------------------------------------------- S5 (V TURNING POINT)
s = slides[4]
clear_slide(s)
base_content(s, 4, "TURNING POINT",
             "Two hidden assumptions, found by first principles", title_size=26)
# left text column
lx, lw = Inches(0.7), Inches(5.95)
add_text(s, lx, Inches(2.3), lw, Inches(2.0),
         [
           [("①  ", {'bold': True, 'color': GOLD, 'size': 16}),
            ("The projection centre had been pinned to the ego origin since day one "
             "— never a design variable. Moving it to the ring-camera centroid cut "
             "projection residuals ", {'size': 13.5, 'color': INK}),
            ("18–96×.", {'size': 13.5, 'color': DARKRED, 'bold': True})],
         ], line_spacing=1.15)
add_text(s, lx, Inches(4.15), lw, Inches(2.6),
         [
           [("②  ", {'bold': True, 'color': GOLD, 'size': 16}),
            ("A fourth error source nobody modeled: the seven cameras fire up to ",
             {'size': 13.5, 'color': INK}),
            ("±22.5 ms apart.", {'size': 13.5, 'color': DARKRED, 'bold': True}),
            (" At 15 m/s a car moves ~0.7 m between adjacent cameras' exposures — "
             "a ", {'size': 13.5, 'color': INK}),
            ("TIME problem no alignment can fix.", {'size': 13.5, 'color': DARKRED, 'bold': True})],
         ], line_spacing=1.15)
# right placeholder D2
PLACE_GEOM['D2'] = (5,) + placeholder(s, "[[D2]]", Inches(6.95), Inches(2.3),
                                      Inches(5.7), Inches(4.0))

# ---------------------------------------------------------------- S6 (VI METHOD I)
s = slides[5]
clear_slide(s)
base_content(s, 5, "METHOD I", "The pipeline: evidence-gated compositing",
             title_size=28)
# full-width placeholder D1
PLACE_GEOM['D1'] = (6,) + placeholder(s, "[[D1]]", Inches(0.7), Inches(2.15),
                                      Inches(11.93), Inches(4.2))
# compact line below
add_text(s, Inches(0.7), Inches(6.45), Inches(11.93), Inches(0.95),
         [
           [("Virtual centre → depth → photometric gain → EMC (ego shutter) "
             "→ segmentation evidence → one body · one camera · one time "
             "(+OMC) → morph + content seam → temporal consensus.   ",
             {'size': 12.5, 'color': INK}),
            ("Rule zero: never average geometry; abstain where evidence is insufficient.",
             {'size': 12.5, 'color': DARKRED, 'bold': True})],
         ], line_spacing=1.1)

# ---------------------------------------------------------------- S7 (VII METHOD II)
s = slides[6]
clear_slide(s)
base_content(s, 6, "METHOD II", "Moving objects: one body, one camera, one time",
             title_size=27)
# image left (tall ~square)
add_image(s, img(r"agent\figs\fig1_porsche_before_after.png"),
          Inches(0.7), Inches(2.25), Inches(5.4), Inches(4.6),
          align='left', valign='top')
# bullets right
add_text(s, Inches(6.5), Inches(2.35), Inches(6.13), Inches(4.4),
         bullets([
           "Boxes give identity, masks give geometry",
           ("The whole car is rendered from ONE camera at ONE instant — seam-cut "
            "cars become impossible by construction", {'bold': True, 'color': DARKRED}),
           "OMC: the object's own shutter displacement measured by ECC (du = +6 px) and compensated",
           "Blend only where views agree (content seam)",
         ], size=14), line_spacing=1.12, space_after=12)

# ---------------------------------------------------------------- S8 (VIII RESULTS)
s = slides[7]
clear_slide(s)
base_content(s, 7,
             "RESULTS",
             "Scene band: source-faithful across 5 scenes, zero per-scene parameters",
             title_size=23)
# wide image top
add_image(s, img(r"agent\figs\fig2_bmw_final_pano.jpg"),
          Inches(0.7), Inches(2.1), Inches(11.93), Inches(2.6),
          align='center', valign='top')
# bullets below
add_text(s, Inches(0.7), Inches(5.0), Inches(11.93), Inches(2.0),
         bullets([
           ("All user-visible defect classes closed on the hardest scene; residuals "
            "adjudicated as real source content", {'bold': True, 'color': DARKRED}),
           "Same code, no tuning: bmw / downtown / crowd / clean / highway",
           "Graceful no-LiDAR degradation: the moving car still renders intact (same OMC du = +6)",
         ], size=14), line_spacing=1.12, space_after=10)

# ---------------------------------------------------------------- S9 (IX GROUND)
s = slides[8]
clear_slide(s)
base_content(s, 8, "COMPLETING THE SPHERE · GROUND",
             "Ground: not generated — reprojected from time", title_size=26)
lx, lw = Inches(0.7), Inches(5.95)
add_text(s, lx, Inches(2.25), lw, Inches(1.7),
         [
           [("The road under the car NOW was fully visible to the cameras seconds "
             "earlier / later → deterministic reprojection of real pixels ",
             {'size': 13.5, 'color': INK}),
            ("(94.5–100% coverage, lane lines continuous through the nadir, ego hood "
             "geometrically removed).", {'size': 13.5, 'color': DARKRED, 'bold': True})],
         ], line_spacing=1.15)
add_text(s, lx, Inches(4.1), lw, Inches(2.7),
         bullets([
           "Candidate frames chosen by geometry over the whole log (a fixed time window fails when the ego idles at a light)",
           "Exact ray test removes the source car's own hood / cabin reflections",
           "Nadir rendered at the evidence's true optical resolution — nothing invented",
         ], size=12.5), line_spacing=1.1, space_after=9)
PLACE_GEOM['D3'] = (9,) + placeholder(s, "[[D3]]", Inches(6.95), Inches(2.3),
                                      Inches(5.7), Inches(4.0))

# ---------------------------------------------------------------- S10 (X SKY)
s = slides[9]
clear_slide(s)
base_content(s, 9, "COMPLETING THE SPHERE · SKY",
             "Sky: the only generated layer", title_size=27)
# I/O block (3 chips) across upper area, full width
io_t = Inches(2.15)
io_h = Inches(0.95)
io_items = [
    ("INPUT", "scene-band panorama + sky mask (the only evidence-free region)"),
    ("MODEL", "FLUX.1-Fill (mask-conditioned inpainting) + DiT360 panorama LoRA (ERP geometry only)"),
    ("OUTPUT", "complete sphere; every pixel outside the mask byte-identical to input"),
]
iw = Inches(3.78)
igap = Inches(0.30)
ix0 = Inches(0.7)
for i, (k, v) in enumerate(io_items):
    cx = ix0 + i * (iw + igap)
    add_rect(s, cx, io_t, iw, io_h, fill=CARD)
    add_rect(s, cx, io_t, Inches(0.06), io_h, fill=GOLD)
    add_text(s, cx + Inches(0.18), io_t + Inches(0.08), iw - Inches(0.3), Inches(0.28),
             k, size=11.5, bold=True, color=DARKRED)
    add_text(s, cx + Inches(0.18), io_t + Inches(0.38), iw - Inches(0.3), Inches(0.55),
             v, size=10.5, color=INK, line_spacing=1.0)
    if i < 2:
        add_text(s, cx + iw, io_t + Inches(0.25), igap, Inches(0.45),
                 "→", size=18, bold=True, color=GRAY_MUTED, align=PP_ALIGN.CENTER)
# bullets (left half) + before/after image (left) + D4 (right)
add_text(s, Inches(0.7), Inches(3.35), Inches(5.55), Inches(1.5),
         bullets([
           ("FLUX.1-Fill is the engine — continuation of observed clouds is its "
            "training objective; DiT360's only role is its panorama LoRA",
            {'bold': True, 'color': DARKRED}),
           "Auto-prompt (sunny / dusk / overcast) chosen from observed sky-band statistics — 5/5 scenes correct",
         ], size=12), line_spacing=1.08, space_after=8)
# before/after image bottom-left
add_image(s, img(r"agent\figs\fig5_complete_before_after.jpg"),
          Inches(0.7), Inches(4.95), Inches(5.55), Inches(2.0),
          align='left', valign='top')
# D4 right half
PLACE_GEOM['D4'] = (10,) + placeholder(s, "[[D4]]", Inches(6.65), Inches(3.35),
                                       Inches(6.0), Inches(3.6))

# ---------------------------------------------------------------- S11 (XI FINAL RESULTS)
s = slides[10]
clear_slide(s)
base_content(s, 10, "FINAL RESULTS", "Complete spheres, multi-weather", title_size=30)
# big tall 3-stack, centered, height-limited to ~4.45 in (max ~6.2 wide), caption below
add_image(s, img(r"agent\figs\fig6_v8_multiweather.jpg"),
          Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.45),
          align='center', valign='top')
add_text(s, Inches(0.7), Inches(6.55), Inches(11.93), Inches(0.6),
         "bmw sunny cumulus · downtown auto-detected dusk · highway altocumulus "
         "continuation — scene band & ground 100% real pixels; sky is the only "
         "generated region.",
         size=11.5, color=GRAY_MUTED, align=PP_ALIGN.CENTER, line_spacing=1.05)

# ---------------------------------------------------------------- S12 (XII CLOSING)
s = slides[11]
clear_slide(s)
base_content(s, 11, "CLOSING", "Principles & next", title_size=30)
# left card: principles
add_text(s, Inches(0.7), Inches(2.3), Inches(5.85), Inches(0.34),
         "PRINCIPLES", size=12, bold=True, color=GRAY_MUTED)
add_rect(s, Inches(0.7), Inches(2.67), Inches(5.85), Inches(0.012), fill=GOLD)
add_text(s, Inches(0.7), Inches(2.85), Inches(5.85), Inches(3.6),
         bullets([
           ("Source-faithful: never invent geometry", {'bold': True, 'color': DARKRED}),
           "Evidence-gated: every pixel has a provenance — real (scene band, ground) vs generated (sky only)",
           "Abstain honestly where evidence ends",
         ], size=14), line_spacing=1.15, space_after=14)
# right card: next
add_text(s, Inches(7.0), Inches(2.3), Inches(5.85), Inches(0.34),
         "NEXT", size=12, bold=True, color=GRAY_MUTED)
add_rect(s, Inches(7.0), Inches(2.67), Inches(5.85), Inches(0.012), fill=GOLD)
add_text(s, Inches(7.0), Inches(2.85), Inches(5.85), Inches(3.6),
         bullets([
           ("Waymo migration — the generalization gate", {'bold': True, 'color': DARKRED}),
           "Centre contract with the world-model first-frame consumer",
           "Dataset scaling: 75-panorama AV2 set already produced",
         ], size=14), line_spacing=1.15, space_after=14)

# ---------------------------------------------------------------- delete surplus slide 13 (index 12)
xml_slides = prs.slides._sldIdLst
sldIds = list(xml_slides)
xml_slides.remove(sldIds[12])

prs.save(OUT)
print("SAVED", OUT)
print("slides:", len(prs.slides))
print("PLACEHOLDERS:")
for k in ('D1','D2','D3','D4'):
    g = PLACE_GEOM[k]
    print(f"  {k}: slide#={g[0]}  L={g[1]} T={g[2]} W={g[3]} H={g[4]} (inches)")
